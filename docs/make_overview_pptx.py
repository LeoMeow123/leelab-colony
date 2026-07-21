#!/usr/bin/env python3
"""One large connected workflow chart for the Lee Lab Colony DB (+ title slide)."""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.oxml.ns import qn

BLUE=RGBColor(0x25,0x63,0xEB); DBLUE=RGBColor(0x1E,0x3A,0x8A); GREEN=RGBColor(0x15,0x80,0x3D)
DGREEN=RGBColor(0x0F,0x5A,0x2E); AMBER=RGBColor(0xB4,0x53,0x09); PURPLE=RGBColor(0x7C,0x3A,0xED)
TEAL=RGBColor(0x0D,0x94,0x88); SLATE=RGBColor(0x33,0x41,0x55); GREY=RGBColor(0x64,0x74,0x8B)
WHITE=RGBColor(0xFF,0xFF,0xFF); INK=RGBColor(0x1A,0x1D,0x26); CONN=RGBColor(0x94,0xA3,0xB8)

prs=Presentation(); prs.slide_width=Inches(13.333); prs.slide_height=Inches(7.5)
BLANK=prs.slide_layouts[6]; SW=prs.slide_width

def slide(): return prs.slides.add_slide(BLANK)

def textbox(sl,s,l,t,w,h,size=18,color=INK,bold=False,align=PP_ALIGN.LEFT,anchor=MSO_ANCHOR.TOP,italic=False):
    tb=sl.shapes.add_textbox(Inches(l),Inches(t),Inches(w),Inches(h)); tf=tb.text_frame
    tf.word_wrap=True; tf.vertical_anchor=anchor; p=tf.paragraphs[0]; p.alignment=align
    r=p.add_run(); r.text=s; f=r.font; f.size=Pt(size); f.bold=bold; f.italic=italic; f.color.rgb=color; f.name="Calibri"
    return tb

# box() records geometry so connectors can snap to edges
def box(sl,l,t,w,h,text,fill=BLUE,font=WHITE,size=12,bold=True,sub=None,line=None):
    shp=sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(l),Inches(t),Inches(w),Inches(h))
    shp.fill.solid(); shp.fill.fore_color.rgb=fill
    if line is None: shp.line.fill.background()
    else: shp.line.color.rgb=line; shp.line.width=Pt(1.25)
    shp.shadow.inherit=False
    tf=shp.text_frame; tf.word_wrap=True; tf.vertical_anchor=MSO_ANCHOR.MIDDLE
    for m in ("margin_left","margin_right","margin_top","margin_bottom"): setattr(tf,m,Inches(0.05))
    p=tf.paragraphs[0]; p.alignment=PP_ALIGN.CENTER
    r=p.add_run(); r.text=text; r.font.size=Pt(size); r.font.bold=bold; r.font.color.rgb=font; r.font.name="Calibri"
    if sub:
        p2=tf.add_paragraph(); p2.alignment=PP_ALIGN.CENTER
        r2=p2.add_run(); r2.text=sub; r2.font.size=Pt(size-2); r2.font.color.rgb=font; r2.font.name="Calibri"
    return {"l":l,"t":t,"w":w,"h":h,"shp":shp}

def R(b): return (b["l"]+b["w"], b["t"]+b["h"]/2)
def L(b): return (b["l"], b["t"]+b["h"]/2)
def T(b): return (b["l"]+b["w"]/2, b["t"])
def B(b): return (b["l"]+b["w"]/2, b["t"]+b["h"])

def connect(sl,p1,p2,color=CONN,width=1.75,two=False,dash=False):
    c=sl.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,Inches(p1[0]),Inches(p1[1]),Inches(p2[0]),Inches(p2[1]))
    c.line.color.rgb=color; c.line.width=Pt(width)
    ln=c.line._get_or_add_ln()
    if dash:
        d=ln.makeelement(qn('a:prstDash'),{'val':'dash'}); ln.append(d)
    te=ln.makeelement(qn('a:tailEnd'),{'type':'triangle','w':'med','len':'med'}); ln.append(te)
    if two:
        he=ln.makeelement(qn('a:headEnd'),{'type':'triangle','w':'med','len':'med'}); ln.append(he)
    return c

def note(sl,s,l,t,w,color=AMBER,size=10):
    return textbox(sl,s,l,t,w,0.3,size=size,color=color,italic=True,align=PP_ALIGN.CENTER)

def speaker(sl,text):
    sl.notes_slide.notes_text_frame.text=text

# ============================================================ Slide 1: title
s=slide()
bg=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,0,0,SW,prs.slide_height); bg.fill.solid()
bg.fill.fore_color.rgb=DBLUE; bg.line.fill.background(); bg.shadow.inherit=False
textbox(s,"🐭 Lee Lab Mouse Colony Database",1,2.6,11.3,1.2,size=40,color=WHITE,bold=True,align=PP_ALIGN.CENTER)
textbox(s,"How it works, end to end — logins, roles, and why different users see different things",
        1,3.9,11.3,0.8,size=20,color=RGBColor(0xC7,0xD6,0xF5),align=PP_ALIGN.CENTER)
speaker(s,
"Our lab's mouse colony database: a single web page backed by a hosted database (Supabase). "
"This one diagram shows how it works end to end, how several people use it at once (each with "
"their own login), and why different people see and can do different things.\n\n"
"Live at leomeow123.github.io/leelab-colony (moving to colony.kuofenleelab.com). "
"No server to maintain and it runs on free tiers.")

# ============================================================ Slide 2: big chart
s=slide()
bar=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,0,0,SW,Inches(0.75)); bar.fill.solid()
bar.fill.fore_color.rgb=DBLUE; bar.line.fill.background(); bar.shadow.inherit=False
tf=bar.text_frame; tf.margin_left=Inches(0.4); tf.vertical_anchor=MSO_ANCHOR.MIDDLE
r=tf.paragraphs[0].add_run(); r.text="Lee Lab Colony DB — one connected workflow"
r.font.size=Pt(22); r.font.bold=True; r.font.color.rgb=WHITE; r.font.name="Calibri"
# legend (4 roles)
textbox(s,"Roles:",9.15,0.22,0.6,0.32,size=11,color=RGBColor(0xC7,0xD6,0xF5),bold=True)
for i,(lbl,col) in enumerate([("PI",PURPLE),("Mgr",BLUE),("Member",SLATE),("Web-maint.",TEAL)]):
    x=9.75+i*0.82
    sq=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,Inches(x),Inches(0.27),Inches(0.15),Inches(0.15))
    sq.fill.solid(); sq.fill.fore_color.rgb=col; sq.line.fill.background(); sq.shadow.inherit=False
    textbox(s,lbl,x+0.18,0.2,0.65,0.3,size=9,color=WHITE)

# ---- TOP SPINE: users -> login -> app <-> supabase -> tables -> backups
textbox(s,"USERS",0.35,0.95,1.5,0.3,size=11,color=GREY,bold=True)
pi =box(s,0.35,1.30,1.55,0.40,"PI",fill=PURPLE,size=12)
mgr=box(s,0.35,1.74,1.55,0.40,"Manager",fill=BLUE,size=12)
mem=box(s,0.35,2.18,1.55,0.40,"Member",fill=SLATE,size=12)
web=box(s,0.35,2.62,1.55,0.40,"Web-maintainer",fill=TEAL,size=10)
login=box(s,2.55,1.72,2.05,1.05,"Log in",fill=DBLUE,size=12,sub="your own password (set on 1st login) · Leo has a master password")
app =box(s,5.05,1.75,2.5,1.0,"Colony web app",fill=BLUE,size=13,sub="index.html · GitHub Pages")
sb  =box(s,10.15,1.20,2.95,1.02,"Supabase",fill=GREEN,size=13,sub="Postgres + RLS + auto REST API")
tbl =box(s,10.15,2.34,2.95,0.60,"mice · app_users · requests · guardian_angels · audit_log",fill=DGREEN,size=9)
bkp =box(s,10.15,3.06,2.95,0.5,"Daily backups (GitHub Actions)",fill=AMBER,size=11)

for u in (pi,mgr,mem,web): connect(s,R(u),L(login))
connect(s,R(login),L(app))
connect(s,R(app),L(sb),two=True); note(s,"reads / writes  (public key + RLS)",7.45,1.42,2.6,color=GREY)
connect(s,B(sb),T(tbl)); connect(s,B(tbl),T(bkp))

# ---- MIDDLE: what anyone can do in the app
textbox(s,"In the app, anyone can:",5.0,2.86,4.0,0.3,size=11,color=GREY,bold=True,italic=True)
add =box(s,2.55,3.20,2.35,0.9,"Add cohort / single mouse",fill=BLUE,size=12,sub="or Import CSV/Excel (auto-match)")
brw =box(s,5.05,3.20,2.4,0.9,"Browse the colony",fill=SLATE,size=12,sub="search · filter · grouped view")
myc =box(s,7.6,3.20,2.4,0.9,"My colony",fill=PURPLE,size=12,sub="your cohorts · edit / delete")
connect(s,B(app),T(add)); connect(s,B(app),T(brw)); connect(s,B(app),T(myc))
note(s,"↑ owner-scoped: you edit / delete only your own mice (managers & Leo: any)",5.0,4.12,5.1,color=AMBER)
# morale extras (right column, under backups)
textbox(s,"Also in the app:",10.15,3.66,3.0,0.3,size=10,color=GREY,bold=True,italic=True)
mor =box(s,10.15,3.96,2.95,0.62,"😇 Guardian Angels · 🕯️ In Memoriam",fill=PURPLE,size=10,sub="pet photos + a warm daily greeting")

# ---- BOTTOM: requests + breeding board + email pipeline
textbox(s,"Requests, the breeding board & email:",0.35,4.58,6.0,0.3,size=11,color=GREY,bold=True,italic=True)
req =box(s,0.35,4.95,2.4,1.0,"File a request",fill=BLUE,size=12,sub="task type · which mice / what's needed")
asg =box(s,2.95,4.95,2.35,1.0,"Assign to people",fill=GREEN,size=12,sub="one or several · picked by role")
act =box(s,5.5,4.95,2.35,1.0,"Done → close",fill=AMBER,size=12,sub="assignee marks done; manager closes")
noti=box(s,8.9,4.70,4.2,0.62,"Notify the assignees",fill=PURPLE,size=12)
send=box(s,8.9,5.48,4.2,0.92,"Edge Function → Resend → email",fill=DBLUE,size=11,sub="with a link back to the app (else draft)")
connect(s,B(mem),(1.55,4.95),color=CONN,dash=True)   # a member starts a request
connect(s,R(req),L(asg)); connect(s,R(asg),L(act))
connect(s,R(asg),L(noti)); connect(s,B(noti),T(send))
note(s,"'Request for mice' / 'Breed' collect on the Mice-needed board",3.0,6.0,4.9,color=AMBER)

# ---- takeaway strip
strip=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.35),Inches(6.78),Inches(12.75),Inches(0.52))
strip.fill.solid(); strip.fill.fore_color.rgb=RGBColor(0xEF,0xF3,0xFB); strip.line.color.rgb=BLUE
strip.line.width=Pt(1); strip.shadow.inherit=False
tf=strip.text_frame; tf.vertical_anchor=MSO_ANCHOR.MIDDLE; p=tf.paragraphs[0]; p.alignment=PP_ALIGN.CENTER
r=p.add_run()
r.text=("Same data, different views: everyone has their own login · 'My colony' + delete/edit are scoped to you · "
        "only Managers/PI/Web-maintainer act on any mouse · every change is logged in audit_log")
r.font.size=Pt(12); r.font.color.rgb=DBLUE; r.font.name="Calibri"; r.font.bold=True

speaker(s,
"WALKTHROUGH (top row left→right, then the bottom lane)\n"
"Everyone signs in with their NAME and their OWN password. The first time, you use a shared default "
"password once and the app then forces you to set a personal one — so no one can log into anyone else's "
"account. Your name is attached to everything you do. The Web-maintainer (Leo) has two backdoors: a "
"'Log in as' button (impersonation, logged) and a master password that works on any account.\n"
"\n"
"HOW SUPABASE WORKS\n"
"- Supabase is a hosted PostgreSQL database that also auto-generates a web (REST) API in front of every "
"table, plus Row-Level Security and small server-side functions.\n"
"- The web page is one static HTML file (GitHub Pages). It has no backend of its own — it calls Supabase's "
"REST API directly, sending a public 'publishable' key with each request.\n"
"- Data lives in a few tables: mice (the hub), app_users (people + role + email + each person's password "
"hash), requests, guardian_angels + angel_reactions (pets), audit_log, app_config. A view (mouse_v) "
"computes each mouse's age and staleness from its date of birth on the fly.\n"
"- Anything that needs a secret (sending email) runs in a Supabase Edge Function server-side, which holds "
"the mail key and sends with a link back to the app.\n"
"- A GitHub Action dumps every table to a private repo once a day (restorable backups) and keeps the free "
"database awake.\n"
"\n"
"THE ROLES (app_users.role)\n"
"- Member: files requests, adds/imports mice, manages their own colony (edit/delete only their own).\n"
"- Manager: all of the above, plus approve/deny requests, transfer ownership, and act on any mouse "
"(the colony/breeding manager, e.g. Bertha).\n"
"- PI: same authority as a manager, for oversight.\n"
"- Web-maintainer (Leo): full access, plus the password/impersonation backdoors.\n"
"\n"
"HOW DIFFERENT USERS SEE (AND CAN DO) DIFFERENT THINGS\n"
"- It is the same shared database, but the app shapes each person's view by OWNERSHIP, ROLE, and who "
"you are logged in as:\n"
"   - 'My colony' shows only mice you own or created; you can group and expand them.\n"
"   - Delete and batch-edit are owner-scoped: a member only touches their own mice (enforced with a "
"server-side owner filter); managers / PI / Web-maintainer can act on any.\n"
"   - Requests are assigned to specific people; approve / close buttons appear only for managers or the "
"assignees; only the assignees are emailed. 'Request for mice' and 'Breed' collect on the Mice-needed board.\n"
"   - The password setters and backdoors are Web-maintainer-only.\n"
"- Honest point for the lab: this is still a SOFT model — one public key, and the password hashes live in "
"a table that key can read, so a determined technical person could bypass the UI. Per-user passwords stop "
"casual cross-login, roles/ownership decide what is shown and who can act, and every change is in the audit "
"log. For hard per-person enforcement we'd switch to Supabase Auth and write those rules into Row-Level "
"Security — with no change to the data.\n"
"\n"
"AND A NICE TOUCH: everyone uploads their pets as 'Guardian Angels', and each visit a random pet greets "
"you with a warm message; departed mice are honored in 'In Memoriam'.")

out="/home/exx/leelab-colony/docs/Colony_DB_Overview.pptx"; prs.save(out)
print("saved",out,"slides:",len(prs.slides._sldIdLst))

# overlap sanity check on boxes
boxes={"pi":pi,"mgr":mgr,"mem":mem,"web":web,"login":login,"app":app,"sb":sb,"tbl":tbl,"bkp":bkp,
       "add":add,"brw":brw,"myc":myc,"mor":mor,"req":req,"asg":asg,"act":act,"noti":noti,"send":send}
def ov(a,b):
    return not (a["l"]+a["w"]<=b["l"] or b["l"]+b["w"]<=a["l"] or a["t"]+a["h"]<=b["t"] or b["t"]+b["h"]<=a["t"])
names=list(boxes); hit=[]
for i in range(len(names)):
    for j in range(i+1,len(names)):
        if ov(boxes[names[i]],boxes[names[j]]): hit.append((names[i],names[j]))
print("overlaps:",hit if hit else "none")
